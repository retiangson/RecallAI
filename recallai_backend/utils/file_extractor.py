import io
import zipfile
from typing import List

import pdfplumber
import docx              # python-docx
from pptx import Presentation   # python-pptx

from openai import OpenAI
client = OpenAI()


# --------------------------------------------------
# LOCAL EXTRACTION (NO GPT)
# --------------------------------------------------
def extract_text_local(file_bytes: bytes, filename: str) -> str:
    ext = filename.lower().split(".")[-1]

    # ----- PDF -----
    if ext == "pdf":
        try:
            doc = pdfplumber.open(stream=file_bytes, filetype="pdf")
            text = []
            for page in doc:
                text.append(page.get_text())
            return "\n".join(text)
        except Exception as e:
            return f"[PDF extraction error] {e}"

    # ----- DOCX -----
    if ext == "docx":
        try:
            doc_obj = docx.Document(io.BytesIO(file_bytes))
            return "\n".join([p.text for p in doc_obj.paragraphs])
        except Exception as e:
            return f"[DOCX extraction error] {e}"

    # ----- PPTX -----
    if ext == "pptx":
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            return f"[PPTX extraction error] {e}"

    # ----- TXT -----
    if ext == "txt":
        return file_bytes.decode("utf-8", errors="ignore")

    # ----- ZIP (TXT files only) -----
    if ext == "zip":
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                texts = []
                for name in z.namelist():
                    if name.lower().endswith(".txt"):
                        texts.append(
                            z.read(name).decode("utf-8", errors="ignore")
                        )
            return "\n".join(texts)
        except Exception as e:
            return f"[ZIP extraction error] {e}"

    # Fallback
    return "Unable to extract text locally."


# --------------------------------------------------
# OPTIONAL GPT CLEANING
# --------------------------------------------------
def clean_text_gpt(raw_text: str) -> str:
    """Optional GPT cleaning. Does NOT summarize. Normalizes only."""
    completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {
                "role": "system",
                "content": (
                    "Normalize this text. "
                    "Do NOT summarize. "
                    "Do NOT remove details. "
                    "Just clean spacing, line breaks, and formatting."
                )
            },
            {
                "role": "user",
                "content": [{"type": "text", "text": raw_text}]
            }
        ]
    )
    return completion.choices[0].message.content


# --------------------------------------------------
# MAIN ENTRY: Combine Local + (optional) GPT cleaning
# --------------------------------------------------
def extract_text_gpt(file_bytes_list: List[bytes], filenames: List[str]) -> str:
    raw = extract_text_local(file_bytes_list[0], filenames[0])

    # OPTIONAL: clean text with GPT
    # (Remove this if you want PURE raw extraction)
    cleaned = clean_text_gpt(raw)

    return cleaned
